#!/usr/bin/env python3
"""Generate SMAP openclaw scenario diagram as a PowerPoint .pptx file."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE_TYPE
import pptx.oxml.ns as nsmap
from lxml import etree

# Slide dimensions: 16:9 (33.87cm x 19.05cm)
SLIDE_WIDTH = Inches(13.33)
SLIDE_HEIGHT = Inches(7.5)

# Slate color palette
C_SLATE_DARK   = RGBColor(0x2E, 0x3D, 0x4A)  # dark slate blue – title bar, headers
C_SLATE_MID    = RGBColor(0x3E, 0x52, 0x64)  # medium slate – column headers
C_SLATE_LIGHT  = RGBColor(0xD0, 0xD8, 0xE0)  # light slate – borders/dividers
C_ORANGE       = RGBColor(0xC8, 0x60, 0x10)  # orange – VM/openclaw/L3 accent
C_ORANGE_LIGHT = RGBColor(0xFF, 0xF0, 0xE0)  # very light orange – VM bg tint
C_BLUE_MID     = RGBColor(0x14, 0x4E, 0x6A)  # steel blue – L2
C_BLUE_LIGHT   = RGBColor(0xE6, 0xF2, 0xF8)  # very light blue – L2 bg tint
C_DRAM_BLUE    = RGBColor(0x2E, 0x52, 0x8A)  # DRAM blue – L1
C_DRAM_LIGHT   = RGBColor(0xE8, 0xEE, 0xF8)  # L1 bg tint
C_WHITE        = RGBColor(0xFF, 0xFF, 0xFF)
C_GRAY_DARK    = RGBColor(0x44, 0x44, 0x44)
C_GRAY_MID     = RGBColor(0x88, 0x88, 0x88)
C_GRAY_LIGHT   = RGBColor(0xF0, 0xF2, 0xF4)


def add_rect(slide, left, top, width, height, fill_rgb=None, line_rgb=None, line_width=Pt(1)):
    """Add a rectangle shape."""
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        left, top, width, height
    )
    if fill_rgb:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_rgb
    else:
        shape.fill.background()
    if line_rgb:
        shape.line.color.rgb = line_rgb
        shape.line.width = line_width
    else:
        shape.line.fill.background()
    return shape


def add_textbox(slide, left, top, width, height, text, font_size=Pt(11),
                font_bold=False, font_color=None, align=PP_ALIGN.LEFT,
                wrap=True, font_name="Arial"):
    """Add a textbox."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = font_size
    run.font.bold = font_bold
    if font_color:
        run.font.color.rgb = font_color
    run.font.name = font_name
    return txBox


def add_shape_text(shape, text, font_size=Pt(11), font_bold=False,
                   font_color=None, align=PP_ALIGN.CENTER, font_name="Arial"):
    """Set text on an existing shape."""
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = font_size
    run.font.bold = font_bold
    if font_color:
        run.font.color.rgb = font_color
    run.font.name = font_name


def add_rounded_rect(slide, left, top, width, height, fill_rgb, line_rgb=None,
                     line_width=Pt(1), corner_size=Emu(91440)):
    """Add a rounded rectangle using MSO_SHAPE_TYPE 5 (ROUNDED_RECTANGLE)."""
    shape = slide.shapes.add_shape(
        5,  # MSO_SHAPE_TYPE.ROUNDED_RECTANGLE
        left, top, width, height
    )
    # Adjust corner radius
    sp = shape.element
    prstGeom = sp.find('.//' + nsmap.qn('a:prstGeom'))
    if prstGeom is not None:
        avLst = prstGeom.find(nsmap.qn('a:avLst'))
        if avLst is not None:
            gd = avLst.find(nsmap.qn('a:gd'))
            if gd is not None:
                gd.set('fmla', 'val 20000')

    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_rgb
    if line_rgb:
        shape.line.color.rgb = line_rgb
        shape.line.width = line_width
    else:
        shape.line.fill.background()
    return shape


def inches(*args):
    return tuple(Inches(a) for a in args)


def create_slide(prs):
    slide_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)
    return slide


def build_pptx():
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT

    slide = create_slide(prs)

    # ── TITLE BAR ──────────────────────────────────────────────────────────────
    title_h = Inches(0.65)
    title_bar = add_rect(slide, 0, 0, SLIDE_WIDTH, title_h, fill_rgb=C_SLATE_DARK)

    # Orange left accent strip
    accent = add_rect(slide, 0, 0, Inches(0.1), title_h, fill_rgb=C_ORANGE)

    # Title text
    add_textbox(
        slide,
        Inches(0.2), Inches(0.08),
        Inches(10), Inches(0.5),
        "openclaw 业务场景  —  SMAP 内存分层优化方案",
        font_size=Pt(18), font_bold=True, font_color=C_WHITE,
        align=PP_ALIGN.LEFT
    )

    # Slide number
    add_textbox(
        slide,
        Inches(12.5), Inches(0.1),
        Inches(0.8), Inches(0.4),
        "01", font_size=Pt(14), font_bold=True, font_color=C_ORANGE,
        align=PP_ALIGN.RIGHT
    )

    # ── CONTENT AREA ───────────────────────────────────────────────────────────
    content_top = Inches(0.75)
    content_h = Inches(6.3)
    col_padding = Inches(0.15)

    # Three main columns
    col1_left = Inches(0.15)
    col1_w = Inches(2.5)
    col2_left = col1_left + col1_w + col_padding
    col2_w = Inches(5.2)
    col3_left = col2_left + col2_w + col_padding
    col3_w = Inches(5.05)

    # ── COLUMN 1: 调用端 ────────────────────────────────────────────────────────
    # Column header
    col1_hdr = add_rect(
        slide, col1_left, content_top, col1_w, Inches(0.38),
        fill_rgb=C_SLATE_MID
    )
    add_shape_text(col1_hdr, "调 用 端", Pt(13), True, C_WHITE)

    # Caller figure box (light background)
    caller_top = content_top + Inches(0.45)
    caller_h = Inches(2.5)
    caller_box = add_rect(
        slide, col1_left, caller_top, col1_w, caller_h,
        fill_rgb=C_GRAY_LIGHT, line_rgb=C_SLATE_LIGHT, line_width=Pt(1)
    )

    # Person icon (simple circle head + body using shapes)
    icon_cx = col1_left + col1_w / 2
    head_size = Inches(0.45)
    head_top = caller_top + Inches(0.2)
    head = slide.shapes.add_shape(
        9,  # OVAL
        icon_cx - head_size / 2, head_top, head_size, head_size
    )
    head.fill.solid()
    head.fill.fore_color.rgb = C_SLATE_MID
    head.line.fill.background()

    # Body (rectangle below head)
    body_top = head_top + head_size + Inches(0.04)
    body_w = Inches(0.32)
    body_h = Inches(0.4)
    body = add_rect(
        slide,
        icon_cx - body_w / 2, body_top,
        body_w, body_h,
        fill_rgb=C_SLATE_MID
    )

    # Label
    add_textbox(
        slide,
        col1_left, body_top + body_h + Inches(0.08),
        col1_w, Inches(0.35),
        "调 用 端", Pt(12), True, C_SLATE_DARK, PP_ALIGN.CENTER
    )

    # Request type chips (2×3 grid)
    request_types = [
        ("网页",  C_SLATE_DARK),
        ("图片",  C_SLATE_DARK),
        ("脚本",  C_ORANGE),
        ("数据",  C_SLATE_DARK),
        ("媒体",  C_ORANGE),
        ("API",   C_SLATE_DARK),
    ]
    chip_w = Inches(1.1)
    chip_h = Inches(0.28)
    chip_gap_x = Inches(0.12)
    chip_gap_y = Inches(0.1)
    chips_start_x = col1_left + Inches(0.08)
    chips_start_y = caller_top + Inches(1.5)

    for i, (label, color) in enumerate(request_types):
        row = i // 2
        col = i % 2
        cx = chips_start_x + col * (chip_w + chip_gap_x)
        cy = chips_start_y + row * (chip_h + chip_gap_y)
        chip = add_rounded_rect(slide, cx, cy, chip_w, chip_h, fill_rgb=color)
        add_shape_text(chip, label, Pt(9), False, C_WHITE)

    # Description text below caller box
    desc_top = caller_top + caller_h + Inches(0.12)
    add_textbox(
        slide,
        col1_left, desc_top, col1_w, Inches(0.6),
        "批量任务 · 采集脚本\n自动化测试 · API 调用",
        Pt(9), False, C_GRAY_DARK, PP_ALIGN.CENTER
    )

    # ── CURVED ARROW col1 -> col2 ───────────────────────────────────────────
    # We'll use a simple connector arrow
    arrow1_left = col1_left + col1_w + Inches(0.02)
    arrow1_top = content_top + Inches(1.5)
    conn = slide.shapes.add_connector(
        1,  # STRAIGHT
        arrow1_left, arrow1_top,
        col2_left - Inches(0.02), arrow1_top
    )
    conn.line.color.rgb = C_SLATE_MID
    conn.line.width = Pt(1.5)

    # Arrow head text (manual)
    add_textbox(
        slide,
        arrow1_left + Inches(0.02), arrow1_top - Inches(0.18),
        Inches(0.12), Inches(0.2),
        "→", Pt(11), False, C_SLATE_MID, PP_ALIGN.CENTER
    )

    # ── COLUMN 2: VM 虚机层 ─────────────────────────────────────────────────────
    col2_hdr = add_rect(
        slide, col2_left, content_top, col2_w, Inches(0.38),
        fill_rgb=C_ORANGE
    )
    add_shape_text(col2_hdr, "宿主机虚机层  ·  openclaw browser", Pt(13), True, C_WHITE)

    vm_area_top = content_top + Inches(0.45)
    vm_area_h = Inches(5.5)
    vm_area = add_rect(
        slide, col2_left, vm_area_top, col2_w, vm_area_h,
        fill_rgb=C_ORANGE_LIGHT, line_rgb=RGBColor(0xE0, 0xB0, 0x80), line_width=Pt(1)
    )

    # VM grid: 2 rows × 3 cols = 6 VMs
    vm_cols = 3
    vm_rows = 2
    vm_w = Inches(1.45)
    vm_h = Inches(1.7)
    vm_gap_x = Inches(0.18)
    vm_gap_y = Inches(0.22)
    vm_grid_left = col2_left + Inches(0.22)
    vm_grid_top = vm_area_top + Inches(0.22)

    for row in range(vm_rows):
        for col in range(vm_cols):
            idx = row * vm_cols + col
            is_new = idx >= 3  # VMs 4-6 are SMAP expansion
            vx = vm_grid_left + col * (vm_w + vm_gap_x)
            vy = vm_grid_top + row * (vm_h + vm_gap_y)

            vm_fill = RGBColor(0xFF, 0xF8, 0xF0) if not is_new else RGBColor(0xFF, 0xED, 0xD5)
            vm_border = C_ORANGE if not is_new else RGBColor(0xD0, 0x70, 0x10)
            border_w = Pt(1.5) if not is_new else Pt(2)

            vm_box = add_rounded_rect(
                slide, vx, vy, vm_w, vm_h, fill_rgb=vm_fill,
                line_rgb=vm_border, line_width=border_w
            )

            # VM screen icon (monitor shape using rectangle)
            screen_x = vx + Inches(0.2)
            screen_y = vy + Inches(0.12)
            screen_w = Inches(1.05)
            screen_h = Inches(0.65)
            screen = add_rect(
                slide, screen_x, screen_y, screen_w, screen_h,
                fill_rgb=C_SLATE_DARK, line_rgb=C_SLATE_MID, line_width=Pt(0.75)
            )
            # Browser address bar inside screen
            bar = add_rect(
                slide, screen_x + Inches(0.05), screen_y + Inches(0.08),
                screen_w - Inches(0.1), Inches(0.12),
                fill_rgb=C_WHITE, line_rgb=None
            )
            # Browser content area
            content_area = add_rect(
                slide, screen_x + Inches(0.05), screen_y + Inches(0.22),
                screen_w - Inches(0.1), Inches(0.35),
                fill_rgb=RGBColor(0x1A, 0x2A, 0x4A), line_rgb=None
            )
            # Small lines to simulate text
            for li in range(3):
                line_rect = add_rect(
                    slide,
                    screen_x + Inches(0.08), screen_y + Inches(0.24 + li * 0.09),
                    screen_w - Inches(0.2), Inches(0.05),
                    fill_rgb=RGBColor(0x4A, 0x9E, 0xFF), line_rgb=None
                )

            # VM label
            vm_label = f"VM-0{idx+1}"
            add_textbox(
                slide, vx, vy + screen_h + Inches(0.15),
                vm_w, Inches(0.22),
                vm_label, Pt(10), True, C_ORANGE if not is_new else RGBColor(0xD0, 0x70, 0x10),
                PP_ALIGN.CENTER
            )
            # openclaw label
            add_textbox(
                slide, vx, vy + screen_h + Inches(0.36),
                vm_w, Inches(0.2),
                "openclaw", Pt(8.5), False, C_GRAY_DARK, PP_ALIGN.CENTER
            )
            # Star for new VMs
            if is_new:
                add_textbox(
                    slide, vx + vm_w - Inches(0.28), vy + Inches(0.06),
                    Inches(0.25), Inches(0.2),
                    "★", Pt(11), True, C_ORANGE, PP_ALIGN.CENTER
                )

    # SMAP status pill
    pill_top = vm_area_top + vm_rows * (vm_h + vm_gap_y) + vm_gap_y + Inches(0.12)
    pill = add_rounded_rect(
        slide,
        col2_left + Inches(0.5), pill_top,
        col2_w - Inches(1.0), Inches(0.35),
        fill_rgb=RGBColor(0x1A, 0x6A, 0x30),
        line_rgb=None
    )
    add_shape_text(pill, "✓  SMAP 已启用  —  冷页面自动卸载至 L2 / L3", Pt(10), True, C_WHITE)

    # Oversubscription comparison
    ovr_top = pill_top + Inches(0.45)
    # Before
    add_textbox(slide, col2_left + Inches(0.2), ovr_top, Inches(2.0), Inches(0.3),
                "无 SMAP  超分比：1×", Pt(9), False, C_GRAY_DARK, PP_ALIGN.LEFT)
    bar_bg1 = add_rect(slide, col2_left + Inches(0.2), ovr_top + Inches(0.3),
                       Inches(2.0), Inches(0.18), fill_rgb=C_SLATE_LIGHT, line_rgb=None)
    bar_fg1 = add_rect(slide, col2_left + Inches(0.2), ovr_top + Inches(0.3),
                       Inches(1.0), Inches(0.18), fill_rgb=C_GRAY_DARK, line_rgb=None)

    # After
    add_textbox(slide, col2_left + Inches(2.7), ovr_top, Inches(2.3), Inches(0.3),
                "启用 SMAP  超分比：↑ 1.8×", Pt(9), False, C_ORANGE, PP_ALIGN.LEFT)
    bar_bg2 = add_rect(slide, col2_left + Inches(2.7), ovr_top + Inches(0.3),
                       Inches(2.0), Inches(0.18), fill_rgb=C_SLATE_LIGHT, line_rgb=None)
    bar_fg2 = add_rect(slide, col2_left + Inches(2.7), ovr_top + Inches(0.3),
                       Inches(1.8), Inches(0.18), fill_rgb=C_ORANGE, line_rgb=None)

    # ── ARROW col2 -> col3 ──────────────────────────────────────────────────
    conn2 = slide.shapes.add_connector(
        1,
        col3_left - Inches(0.12), content_top + Inches(1.5),
        col3_left - Inches(0.02), content_top + Inches(1.5)
    )
    conn2.line.color.rgb = C_SLATE_MID
    conn2.line.width = Pt(1.5)

    # ── COLUMN 3: 内存分层 ──────────────────────────────────────────────────────
    col3_hdr = add_rect(
        slide, col3_left, content_top, col3_w, Inches(0.38),
        fill_rgb=C_SLATE_MID
    )
    add_shape_text(col3_hdr, "SMAP 三级内存分层", Pt(13), True, C_WHITE)

    mem_top = content_top + Inches(0.45)

    # Vertical connector line between memory tiers
    conn_x = col3_left + col3_w / 2

    # ── L1 DRAM ──────────────────────────────────────────────────────────────
    l1_h = Inches(1.75)
    l1_box = add_rect(
        slide, col3_left, mem_top, col3_w, l1_h,
        fill_rgb=C_DRAM_LIGHT, line_rgb=C_DRAM_BLUE, line_width=Pt(1.5)
    )

    # L1 header
    l1_hdr = add_rect(slide, col3_left, mem_top, col3_w, Inches(0.3),
                      fill_rgb=C_DRAM_BLUE)
    add_shape_text(l1_hdr, "L1  ·  本地 DRAM", Pt(11), True, C_WHITE)

    # RAM stick icons (3 sticks)
    ram_y = mem_top + Inches(0.38)
    ram_h = Inches(0.9)
    ram_w = Inches(0.32)
    ram_gap = Inches(0.12)
    ram_total_w = 3 * ram_w + 2 * ram_gap
    ram_start_x = col3_left + (col3_w - ram_total_w) / 2 - Inches(0.8)

    for i in range(3):
        rx = ram_start_x + i * (ram_w + ram_gap)
        stick = add_rect(slide, rx, ram_y, ram_w, ram_h,
                         fill_rgb=C_DRAM_BLUE, line_rgb=RGBColor(0x1A, 0x38, 0x6A), line_width=Pt(0.75))
        # Chip notch
        notch = add_rect(slide, rx + Inches(0.08), ram_y - Inches(0.06),
                         ram_w - Inches(0.16), Inches(0.08),
                         fill_rgb=C_DRAM_LIGHT, line_rgb=None)
        # Chips on stick (3 small squares)
        for c in range(3):
            chip_y = ram_y + Inches(0.12 + c * 0.22)
            chip_sq = add_rect(slide, rx + Inches(0.06), chip_y,
                               ram_w - Inches(0.12), Inches(0.12),
                               fill_rgb=RGBColor(0x4A, 0x80, 0xC0), line_rgb=None)

    # L1 info text
    add_textbox(slide, col3_left + Inches(1.6), ram_y, Inches(3.3), Inches(0.9),
                "热数据  |  直接访问\n延迟  ~100 ns\n容量  2×4 GiB / VM",
                Pt(9.5), False, C_SLATE_DARK, PP_ALIGN.LEFT)

    # Downward arrow L1->L2
    arr_y1 = mem_top + l1_h
    arr_txt1 = add_rect(slide, col3_left + col3_w/2 - Inches(0.4), arr_y1,
                        Inches(0.8), Inches(0.28), fill_rgb=C_DRAM_BLUE)
    add_shape_text(arr_txt1, "Cache Miss  ↓", Pt(8), False, C_WHITE)

    # ── L2 UB借用内存 ────────────────────────────────────────────────────────────
    l2_top = mem_top + l1_h + Inches(0.28)
    l2_h = Inches(1.75)
    l2_box = add_rect(
        slide, col3_left, l2_top, col3_w, l2_h,
        fill_rgb=C_BLUE_LIGHT, line_rgb=C_BLUE_MID, line_width=Pt(1.5)
    )
    l2_hdr = add_rect(slide, col3_left, l2_top, col3_w, Inches(0.3),
                      fill_rgb=C_BLUE_MID)
    add_shape_text(l2_hdr, "L2  ·  UB 借用内存（远端 DRAM）", Pt(11), True, C_WHITE)

    # Network node icon for L2
    node_y = l2_top + Inches(0.4)
    node_size = Inches(0.28)
    node_gap = Inches(0.5)
    node_cx1 = col3_left + Inches(0.6)
    node_cx2 = col3_left + Inches(1.3)
    node_cy = l2_top + Inches(0.85)

    # Two server nodes
    for ncx in [node_cx1, node_cx2]:
        n = slide.shapes.add_shape(9, ncx - node_size/2, node_cy - node_size/2,
                                   node_size, node_size)
        n.fill.solid()
        n.fill.fore_color.rgb = C_BLUE_MID
        n.line.fill.background()

    # Connecting line between nodes
    conn_l2 = slide.shapes.add_connector(
        1,
        node_cx1, node_cy,
        node_cx2, node_cy
    )
    conn_l2.line.color.rgb = C_BLUE_MID
    conn_l2.line.width = Pt(1.5)

    # L2 info text
    add_textbox(slide, col3_left + Inches(1.6), l2_top + Inches(0.38), Inches(3.3), Inches(0.9),
                "温数据  |  RDMA 访问\n延迟  ~1 μs\nopenEuler UB Stack 驱动",
                Pt(9.5), False, C_SLATE_DARK, PP_ALIGN.LEFT)

    # Downward arrow L2->L3
    arr_y2 = l2_top + l2_h
    arr_txt2 = add_rect(slide, col3_left + col3_w/2 - Inches(0.4), arr_y2,
                        Inches(0.8), Inches(0.28), fill_rgb=C_BLUE_MID)
    add_shape_text(arr_txt2, "页面降冷  ↓", Pt(8), False, C_WHITE)

    # ── L3 NVMe ───────────────────────────────────────────────────────────────
    l3_top = l2_top + l2_h + Inches(0.28)
    l3_h = Inches(1.75)
    l3_box = add_rect(
        slide, col3_left, l3_top, col3_w, l3_h,
        fill_rgb=RGBColor(0xFF, 0xF0, 0xE8), line_rgb=C_ORANGE, line_width=Pt(1.5)
    )
    l3_hdr = add_rect(slide, col3_left, l3_top, col3_w, Inches(0.3),
                      fill_rgb=C_ORANGE)
    add_shape_text(l3_hdr, "L3  ·  NVMe SSD  ★ 核心超分扩展", Pt(11), True, C_WHITE)

    # NVMe drive icon
    drive_x = col3_left + Inches(0.2)
    drive_y = l3_top + Inches(0.4)
    drive_w = Inches(1.1)
    drive_h = Inches(1.1)
    drive = add_rect(slide, drive_x, drive_y, drive_w, drive_h,
                     fill_rgb=C_SLATE_DARK, line_rgb=C_ORANGE, line_width=Pt(1))
    # PCB traces (horizontal lines on drive)
    for li in range(4):
        trace = add_rect(slide, drive_x + Inches(0.08),
                         drive_y + Inches(0.12 + li * 0.2),
                         drive_w - Inches(0.16), Inches(0.08),
                         fill_rgb=C_ORANGE, line_rgb=None)
    # M.2 connector
    conn_strip = add_rect(slide, drive_x + drive_w - Inches(0.18),
                           drive_y + drive_h/2 - Inches(0.25),
                           Inches(0.12), Inches(0.5),
                           fill_rgb=C_ORANGE, line_rgb=None)

    # L3 info text
    add_textbox(slide, col3_left + Inches(1.6), l3_top + Inches(0.38), Inches(3.3), Inches(1.1),
                "冷数据  |  块设备访问\n延迟  ~100 μs\n大幅扩展可用内存容量",
                Pt(9.5), False, C_SLATE_DARK, PP_ALIGN.LEFT)

    # ── FOOTER ─────────────────────────────────────────────────────────────────
    footer_top = SLIDE_HEIGHT - Inches(0.38)
    footer = add_rect(slide, 0, footer_top, SLIDE_WIDTH, Inches(0.38),
                      fill_rgb=C_SLATE_DARK)
    # Orange bottom border
    border = add_rect(slide, 0, footer_top, SLIDE_WIDTH, Inches(0.04),
                      fill_rgb=C_ORANGE)
    add_textbox(
        slide, Inches(0.2), footer_top + Inches(0.04),
        Inches(8), Inches(0.3),
        "字节跳动  火山引擎  AI 服务平台  ·  openEuler OLK-6.6  ·  UB Stack  ·  SMAP 分层方案",
        Pt(8.5), False, C_SLATE_LIGHT, PP_ALIGN.LEFT
    )
    add_textbox(
        slide, Inches(9.5), footer_top + Inches(0.04),
        Inches(3.6), Inches(0.3),
        "Confidential  —  2026", Pt(8.5), False, C_SLATE_LIGHT, PP_ALIGN.RIGHT
    )

    # ── SAVE ───────────────────────────────────────────────────────────────────
    import os
    os.makedirs("diagrams", exist_ok=True)
    out_path = "diagrams/smap_openclaw_scenario.pptx"
    prs.save(out_path)
    print(f"Saved: {out_path}")
    return out_path


if __name__ == "__main__":
    path = build_pptx()
    print(f"Done: {path}")
