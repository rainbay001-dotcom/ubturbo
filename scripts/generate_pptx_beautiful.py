#!/usr/bin/env python3
"""
Generate 10 beautiful PPTX slides for the SMAP openclaw scenario.

Requirements:
    pip install python-pptx Pillow cairosvg

Usage:
    python3 scripts/generate_pptx_beautiful.py

Output:
    diagrams/ppt_beautiful_s01_minimalist.pptx  ... s10_monochrome.pptx
"""

import os
import io
import sys
from pathlib import Path

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    from pptx.util import Inches, Pt
    from pptx.enum.dml import MSO_THEME_COLOR
    import pptx.oxml.ns as nsmap
    from lxml import etree
except ImportError:
    print("ERROR: python-pptx not installed. Run: pip install python-pptx")
    sys.exit(1)

try:
    import cairosvg
    HAS_CAIROSVG = True
except ImportError:
    HAS_CAIROSVG = False
    print("WARNING: cairosvg not installed. SVG icons will be skipped.")
    print("  Install with: pip install cairosvg")

try:
    from PIL import Image
    HAS_PIL = True
except ImportError:
    HAS_PIL = False

# ---------------------------------------------------------------------------
# Constants
# ---------------------------------------------------------------------------
SCRIPT_DIR = Path(__file__).parent
REPO_ROOT = SCRIPT_DIR.parent
ASSETS_DIR = REPO_ROOT / "assets" / "icons"
OUTPUT_DIR = REPO_ROOT / "diagrams"
OUTPUT_DIR.mkdir(exist_ok=True)

# Slide size: 16:9 widescreen (13.33" x 7.5")
SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)

# ---------------------------------------------------------------------------
# Helper: RGB shorthand
# ---------------------------------------------------------------------------
def rgb(r, g, b):
    return RGBColor(r, g, b)

# ---------------------------------------------------------------------------
# Helper: convert SVG file to PNG bytes (via cairosvg)
# ---------------------------------------------------------------------------
def svg_to_png_bytes(svg_path: Path, width: int = 200, height: int = 200) -> bytes | None:
    if not HAS_CAIROSVG:
        return None
    if not svg_path.exists():
        return None
    try:
        return cairosvg.svg2png(
            url=str(svg_path),
            output_width=width,
            output_height=height,
        )
    except Exception as e:
        print(f"  WARNING: cairosvg failed for {svg_path.name}: {e}")
        return None

# ---------------------------------------------------------------------------
# Helper: add a slide to a presentation
# ---------------------------------------------------------------------------
def blank_slide(prs: Presentation):
    layout = prs.slide_layouts[6]  # blank
    return prs.slides.add_slide(layout)

# ---------------------------------------------------------------------------
# Helper: add a filled rectangle
# ---------------------------------------------------------------------------
def add_rect(slide, left, top, width, height,
             fill_color=None, line_color=None, line_width=Pt(1)):
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        left, top, width, height
    )
    fill = shape.fill
    if fill_color:
        fill.solid()
        fill.fore_color.rgb = fill_color
    else:
        fill.background()

    line = shape.line
    if line_color:
        line.color.rgb = line_color
        line.width = line_width
    else:
        line.fill.background()
    return shape

# ---------------------------------------------------------------------------
# Helper: add text box
# ---------------------------------------------------------------------------
def add_textbox(slide, left, top, width, height, text,
                font_size=Pt(12), bold=False, color=None,
                align=PP_ALIGN.LEFT, font_name="Arial"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = font_size
    run.font.bold = bold
    run.font.name = font_name
    if color:
        run.font.color.rgb = color
    return txBox

# ---------------------------------------------------------------------------
# Helper: add image from bytes
# ---------------------------------------------------------------------------
def add_image_bytes(slide, png_bytes: bytes, left, top, width, height):
    if png_bytes is None:
        return None
    img_stream = io.BytesIO(png_bytes)
    return slide.shapes.add_picture(img_stream, left, top, width, height)

# ---------------------------------------------------------------------------
# Helper: add rounded rectangle (via freeform or shape type 5)
# pptx shape type 5 = ROUNDED_RECTANGLE
# ---------------------------------------------------------------------------
def add_rounded_rect(slide, left, top, width, height,
                     fill_color=None, line_color=None, line_width=Pt(1),
                     corner_size=0.1):
    from pptx.util import Emu
    shape = slide.shapes.add_shape(
        5,  # rounded rectangle
        left, top, width, height
    )
    # Set corner radius via XML adjustment
    sp_pr = shape.element.spPr
    prstGeom = sp_pr.find('.//{http://schemas.openxmlformats.org/drawingml/2006/main}prstGeom')
    if prstGeom is not None:
        avLst = prstGeom.find('{http://schemas.openxmlformats.org/drawingml/2006/main}avLst')
        if avLst is not None:
            for gd in avLst.findall('{http://schemas.openxmlformats.org/drawingml/2006/main}gd'):
                if gd.get('name') == 'adj':
                    gd.set('fmla', f'val {int(corner_size * 50000)}')

    fill = shape.fill
    if fill_color:
        fill.solid()
        fill.fore_color.rgb = fill_color
    else:
        fill.background()

    line = shape.line
    if line_color:
        line.color.rgb = line_color
        line.width = line_width
    else:
        line.fill.background()
    return shape

# ---------------------------------------------------------------------------
# Shared content data
# ---------------------------------------------------------------------------
REQUEST_TYPES = [
    ("网页", "🌐"),
    ("图片", "🖼"),
    ("脚本", "📜"),
    ("数据", "📊"),
    ("媒体", "🎬"),
    ("API",  "🔗"),
]

MEMORY_LAYERS = [
    {
        "name": "L1  DRAM",
        "sub": "本地内存",
        "latency": "~100 ns",
        "data": "热数据",
        "icon": "dram_icon.svg",
    },
    {
        "name": "L2  UB借用",
        "sub": "跨机借用内存",
        "latency": "~1 μs",
        "data": "温数据",
        "icon": "ub_memory_icon.svg",
    },
    {
        "name": "L3  NVMe",
        "sub": "本地SSD",
        "latency": "~100 μs",
        "data": "冷数据",
        "icon": "nvme_icon.svg",
    },
]

VMS_EXISTING = ["VM-01", "VM-02", "VM-03"]
VMS_EXTRA    = ["VM-04★", "VM-05★", "VM-06★"]

# ---------------------------------------------------------------------------
# THEME DEFINITIONS  (10 themes)
# Each theme defines a color palette and minor layout tweaks.
# ---------------------------------------------------------------------------
THEMES = [
    # 01 — Clean Minimal (Office Default Blue/Orange)
    {
        "name": "s01_clean_minimal",
        "label": "简约蓝橙",
        "header_bg":    rgb(68, 114, 196),   # #4472C4
        "header_text":  rgb(255, 255, 255),
        "accent":       rgb(197, 90, 17),     # #C55A11
        "card_bg":      rgb(242, 242, 242),
        "card_line":    rgb(68, 114, 196),
        "body_bg":      rgb(255, 255, 255),
        "title_text":   rgb(31, 56, 100),
        "body_text":    rgb(50, 50, 50),
        "vm_fill":      rgb(68, 114, 196),
        "vm_extra":     rgb(197, 90, 17),
        "l1_color":     rgb(68, 114, 196),
        "l2_color":     rgb(0, 112, 192),
        "l3_color":     rgb(197, 90, 17),
        "footer_bg":    rgb(31, 56, 100),
        "footer_text":  rgb(200, 210, 230),
        "arrow_color":  rgb(68, 114, 196),
    },
    # 02 — Corporate Slate (stone blue)
    {
        "name": "s02_corporate_slate",
        "label": "商务石板蓝",
        "header_bg":    rgb(46, 61, 74),
        "header_text":  rgb(255, 255, 255),
        "accent":       rgb(200, 98, 40),
        "card_bg":      rgb(245, 247, 250),
        "card_line":    rgb(46, 61, 74),
        "body_bg":      rgb(255, 255, 255),
        "title_text":   rgb(30, 40, 55),
        "body_text":    rgb(55, 65, 80),
        "vm_fill":      rgb(46, 61, 74),
        "vm_extra":     rgb(200, 98, 40),
        "l1_color":     rgb(46, 61, 74),
        "l2_color":     rgb(30, 90, 130),
        "l3_color":     rgb(180, 80, 20),
        "footer_bg":    rgb(30, 40, 55),
        "footer_text":  rgb(180, 195, 215),
        "arrow_color":  rgb(100, 130, 160),
    },
    # 03 — Teal Fresh
    {
        "name": "s03_teal_fresh",
        "label": "清新青蓝",
        "header_bg":    rgb(0, 128, 128),
        "header_text":  rgb(255, 255, 255),
        "accent":       rgb(230, 126, 34),
        "card_bg":      rgb(240, 252, 252),
        "card_line":    rgb(0, 128, 128),
        "body_bg":      rgb(255, 255, 255),
        "title_text":   rgb(0, 80, 80),
        "body_text":    rgb(40, 60, 60),
        "vm_fill":      rgb(0, 128, 128),
        "vm_extra":     rgb(230, 126, 34),
        "l1_color":     rgb(0, 128, 128),
        "l2_color":     rgb(0, 160, 120),
        "l3_color":     rgb(210, 90, 20),
        "footer_bg":    rgb(0, 80, 80),
        "footer_text":  rgb(180, 230, 230),
        "arrow_color":  rgb(0, 128, 128),
    },
    # 04 — Navy Gold (Executive)
    {
        "name": "s04_navy_gold",
        "label": "深蓝金色",
        "header_bg":    rgb(0, 40, 85),
        "header_text":  rgb(255, 215, 0),
        "accent":       rgb(180, 140, 0),
        "card_bg":      rgb(248, 246, 235),
        "card_line":    rgb(0, 40, 85),
        "body_bg":      rgb(255, 255, 255),
        "title_text":   rgb(0, 40, 85),
        "body_text":    rgb(40, 45, 55),
        "vm_fill":      rgb(0, 40, 85),
        "vm_extra":     rgb(180, 140, 0),
        "l1_color":     rgb(0, 40, 85),
        "l2_color":     rgb(10, 80, 150),
        "l3_color":     rgb(160, 120, 0),
        "footer_bg":    rgb(0, 30, 65),
        "footer_text":  rgb(220, 200, 150),
        "arrow_color":  rgb(0, 80, 160),
    },
    # 05 — Steel Blue (calm professional)
    {
        "name": "s05_steel_blue",
        "label": "钢蓝专业",
        "header_bg":    rgb(70, 110, 160),
        "header_text":  rgb(255, 255, 255),
        "accent":       rgb(185, 85, 30),
        "card_bg":      rgb(240, 244, 250),
        "card_line":    rgb(70, 110, 160),
        "body_bg":      rgb(255, 255, 255),
        "title_text":   rgb(40, 70, 110),
        "body_text":    rgb(50, 60, 80),
        "vm_fill":      rgb(70, 110, 160),
        "vm_extra":     rgb(185, 85, 30),
        "l1_color":     rgb(70, 110, 160),
        "l2_color":     rgb(40, 130, 180),
        "l3_color":     rgb(185, 85, 30),
        "footer_bg":    rgb(35, 55, 90),
        "footer_text":  rgb(190, 210, 235),
        "arrow_color":  rgb(70, 130, 190),
    },
    # 06 — Forest Green
    {
        "name": "s06_forest_green",
        "label": "森绿商务",
        "header_bg":    rgb(34, 85, 60),
        "header_text":  rgb(255, 255, 255),
        "accent":       rgb(210, 100, 20),
        "card_bg":      rgb(238, 248, 242),
        "card_line":    rgb(34, 85, 60),
        "body_bg":      rgb(255, 255, 255),
        "title_text":   rgb(20, 65, 40),
        "body_text":    rgb(40, 60, 50),
        "vm_fill":      rgb(34, 85, 60),
        "vm_extra":     rgb(210, 100, 20),
        "l1_color":     rgb(34, 85, 60),
        "l2_color":     rgb(0, 140, 80),
        "l3_color":     rgb(180, 80, 10),
        "footer_bg":    rgb(20, 55, 38),
        "footer_text":  rgb(180, 225, 200),
        "arrow_color":  rgb(0, 140, 80),
    },
    # 07 — Purple Tech
    {
        "name": "s07_purple_tech",
        "label": "科技紫",
        "header_bg":    rgb(75, 0, 130),
        "header_text":  rgb(255, 255, 255),
        "accent":       rgb(220, 60, 120),
        "card_bg":      rgb(248, 240, 255),
        "card_line":    rgb(120, 60, 180),
        "body_bg":      rgb(255, 255, 255),
        "title_text":   rgb(60, 0, 100),
        "body_text":    rgb(50, 30, 70),
        "vm_fill":      rgb(100, 40, 160),
        "vm_extra":     rgb(220, 60, 120),
        "l1_color":     rgb(100, 40, 160),
        "l2_color":     rgb(60, 100, 200),
        "l3_color":     rgb(200, 40, 100),
        "footer_bg":    rgb(50, 0, 90),
        "footer_text":  rgb(210, 180, 240),
        "arrow_color":  rgb(120, 60, 180),
    },
    # 08 — Warm Terracotta
    {
        "name": "s08_warm_terracotta",
        "label": "暖色陶土",
        "header_bg":    rgb(155, 65, 40),
        "header_text":  rgb(255, 255, 255),
        "accent":       rgb(55, 100, 160),
        "card_bg":      rgb(255, 245, 240),
        "card_line":    rgb(155, 65, 40),
        "body_bg":      rgb(255, 255, 255),
        "title_text":   rgb(120, 45, 25),
        "body_text":    rgb(60, 45, 40),
        "vm_fill":      rgb(155, 65, 40),
        "vm_extra":     rgb(55, 100, 160),
        "l1_color":     rgb(155, 65, 40),
        "l2_color":     rgb(80, 140, 190),
        "l3_color":     rgb(120, 45, 25),
        "footer_bg":    rgb(100, 40, 20),
        "footer_text":  rgb(240, 210, 200),
        "arrow_color":  rgb(155, 65, 40),
    },
    # 09 — Soft Lavender (SaaS style)
    {
        "name": "s09_soft_lavender",
        "label": "柔和薰衣草",
        "header_bg":    rgb(100, 80, 170),
        "header_text":  rgb(255, 255, 255),
        "accent":       rgb(240, 130, 50),
        "card_bg":      rgb(245, 243, 255),
        "card_line":    rgb(150, 130, 210),
        "body_bg":      rgb(255, 255, 255),
        "title_text":   rgb(70, 50, 140),
        "body_text":    rgb(60, 50, 80),
        "vm_fill":      rgb(100, 80, 170),
        "vm_extra":     rgb(240, 130, 50),
        "l1_color":     rgb(100, 80, 170),
        "l2_color":     rgb(60, 130, 200),
        "l3_color":     rgb(210, 100, 30),
        "footer_bg":    rgb(70, 50, 130),
        "footer_text":  rgb(220, 210, 250),
        "arrow_color":  rgb(130, 100, 200),
    },
    # 10 — Monochrome (only openclaw orange pops)
    {
        "name": "s10_monochrome",
        "label": "单色高级灰",
        "header_bg":    rgb(40, 40, 40),
        "header_text":  rgb(255, 255, 255),
        "accent":       rgb(234, 88, 12),    # only orange pops
        "card_bg":      rgb(246, 246, 246),
        "card_line":    rgb(160, 160, 160),
        "body_bg":      rgb(255, 255, 255),
        "title_text":   rgb(25, 25, 25),
        "body_text":    rgb(80, 80, 80),
        "vm_fill":      rgb(90, 90, 90),
        "vm_extra":     rgb(234, 88, 12),   # orange = new VMs
        "l1_color":     rgb(100, 100, 100),
        "l2_color":     rgb(130, 130, 130),
        "l3_color":     rgb(234, 88, 12),   # orange = L3 highlight
        "footer_bg":    rgb(30, 30, 30),
        "footer_text":  rgb(180, 180, 180),
        "arrow_color":  rgb(140, 140, 140),
    },
]

# ---------------------------------------------------------------------------
# Core slide builder
# ---------------------------------------------------------------------------
def build_slide(theme: dict) -> Presentation:
    """Build a single-slide PPTX using the given theme."""
    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H
    slide = blank_slide(prs)

    # ------------------------------------------------------------------
    # Layout dimensions (inches, converted via Inches() when needed)
    # ------------------------------------------------------------------
    HEADER_H    = 0.75
    FOOTER_H    = 0.42
    BODY_TOP    = HEADER_H
    BODY_H      = 7.5 - HEADER_H - FOOTER_H
    COL_W       = 13.33 / 3          # 3 equal columns
    PAD         = 0.18               # inner padding
    TITLE_LINE_H= 0.07               # thin accent line below header

    # ------------------------------------------------------------------
    # 1. Background
    # ------------------------------------------------------------------
    add_rect(slide,
             Inches(0), Inches(0), SLIDE_W, SLIDE_H,
             fill_color=theme["body_bg"])

    # ------------------------------------------------------------------
    # 2. Header bar
    # ------------------------------------------------------------------
    add_rect(slide,
             Inches(0), Inches(0), SLIDE_W, Inches(HEADER_H),
             fill_color=theme["header_bg"])

    # Orange accent line on header left
    add_rect(slide,
             Inches(0), Inches(0), Inches(0.06), Inches(HEADER_H),
             fill_color=theme["accent"])

    # Header title
    add_textbox(slide,
                Inches(0.16), Inches(0.06), Inches(9), Inches(HEADER_H - 0.06),
                "SMAP 内存分层超分方案  ·  openclaw 业务场景",
                font_size=Pt(17), bold=True,
                color=theme["header_text"],
                align=PP_ALIGN.LEFT, font_name="Arial")

    # Header subtitle (right side)
    add_textbox(slide,
                Inches(9.2), Inches(0.06), Inches(4.0), Inches(HEADER_H - 0.06),
                "字节跳动火山 AI 服务平台",
                font_size=Pt(10.5),
                color=rgb(210, 220, 240) if theme["header_bg"][0] < 100
                      else rgb(240, 240, 240),
                align=PP_ALIGN.RIGHT, font_name="Arial")

    # Thin separator line
    add_rect(slide,
             Inches(0), Inches(HEADER_H), SLIDE_W, Inches(TITLE_LINE_H),
             fill_color=theme["accent"])

    # ------------------------------------------------------------------
    # 3. Column backgrounds (light card per column)
    # ------------------------------------------------------------------
    col_positions = [0, COL_W, COL_W * 2]
    col_labels    = ["调 用 端", "虚机层  ·  openclaw", "SMAP 内存分层"]

    for i, (cx, cl) in enumerate(zip(col_positions, col_labels)):
        # Card bg
        add_rect(slide,
                 Inches(cx + PAD * 0.5), Inches(BODY_TOP + 0.05),
                 Inches(COL_W - PAD), Inches(BODY_H - 0.10),
                 fill_color=theme["card_bg"],
                 line_color=theme["card_line"],
                 line_width=Pt(0.75))

        # Column label
        add_textbox(slide,
                    Inches(cx + PAD), Inches(BODY_TOP + 0.10),
                    Inches(COL_W - PAD * 2), Inches(0.30),
                    cl,
                    font_size=Pt(11), bold=True,
                    color=theme["title_text"],
                    align=PP_ALIGN.CENTER, font_name="Arial")

        # Underline for column label
        add_rect(slide,
                 Inches(cx + PAD + 0.3), Inches(BODY_TOP + 0.40),
                 Inches(COL_W - PAD * 2 - 0.6), Inches(0.025),
                 fill_color=theme["accent"] if i == 1 else theme["card_line"])

    # ------------------------------------------------------------------
    # 4. COLUMN 1 — Caller section
    # ------------------------------------------------------------------
    C1_LEFT = col_positions[0]
    ICON_Y  = BODY_TOP + 0.50

    # Caller icon (SVG → PNG)
    caller_png = svg_to_png_bytes(ASSETS_DIR / "caller_icon.svg", 160, 160)
    if caller_png:
        add_image_bytes(slide, caller_png,
                        Inches(C1_LEFT + 1.05), Inches(ICON_Y),
                        Inches(2.2), Inches(2.2))
    else:
        # Fallback: simple rectangle placeholder
        add_rounded_rect(slide,
                         Inches(C1_LEFT + 1.0), Inches(ICON_Y),
                         Inches(2.3), Inches(2.3),
                         fill_color=theme["l1_color"], corner_size=0.2)
        add_textbox(slide,
                    Inches(C1_LEFT + 1.0), Inches(ICON_Y + 0.8),
                    Inches(2.3), Inches(0.5), "🖥 调用端",
                    font_size=Pt(14), bold=True,
                    color=rgb(255,255,255), align=PP_ALIGN.CENTER)

    add_textbox(slide,
                Inches(C1_LEFT + PAD), Inches(ICON_Y + 2.2),
                Inches(COL_W - PAD * 2), Inches(0.32),
                "调 用 端",
                font_size=Pt(13), bold=True,
                color=theme["title_text"], align=PP_ALIGN.CENTER)
    add_textbox(slide,
                Inches(C1_LEFT + PAD), Inches(ICON_Y + 2.52),
                Inches(COL_W - PAD * 2), Inches(0.28),
                "批量任务 / API调用 / 自动化测试",
                font_size=Pt(9),
                color=theme["body_text"], align=PP_ALIGN.CENTER)

    # Request type chips (2 rows × 3)
    chip_types = [("🌐 网页", 0), ("🖼 图片", 1), ("📜 脚本", 2),
                  ("📊 数据", 3), ("🎬 媒体", 4), ("🔗 API",  5)]
    CHIP_W  = 1.05
    CHIP_H  = 0.30
    CHIP_ROW_Y = ICON_Y + 2.95
    CHIP_GAP_X = 0.06
    chip_x_start = C1_LEFT + PAD + 0.06
    for idx, (label, _) in enumerate(chip_types):
        row = idx // 3
        col = idx % 3
        cx = chip_x_start + col * (CHIP_W + CHIP_GAP_X)
        cy = CHIP_ROW_Y + row * (CHIP_H + 0.06)
        add_rounded_rect(slide,
                         Inches(cx), Inches(cy),
                         Inches(CHIP_W), Inches(CHIP_H),
                         fill_color=theme["l1_color"] if idx % 2 == 0
                                    else theme["accent"],
                         corner_size=0.4)
        add_textbox(slide,
                    Inches(cx), Inches(cy),
                    Inches(CHIP_W), Inches(CHIP_H),
                    label,
                    font_size=Pt(8.5), bold=False,
                    color=rgb(255, 255, 255),
                    align=PP_ALIGN.CENTER)

    # Arrow from col 1 → col 2
    add_textbox(slide,
                Inches(C1_LEFT + COL_W - 0.40), Inches(BODY_TOP + BODY_H / 2 - 0.15),
                Inches(0.55), Inches(0.35),
                "→",
                font_size=Pt(22), bold=True,
                color=theme["arrow_color"], align=PP_ALIGN.CENTER)

    # ------------------------------------------------------------------
    # 5. COLUMN 2 — VM / openclaw section
    # ------------------------------------------------------------------
    C2_LEFT = col_positions[1]

    # openclaw icon
    openclaw_png = svg_to_png_bytes(ASSETS_DIR / "openclaw_icon.svg", 140, 140)
    if openclaw_png:
        add_image_bytes(slide, openclaw_png,
                        Inches(C2_LEFT + 1.35), Inches(BODY_TOP + 0.50),
                        Inches(1.75), Inches(1.75))
    else:
        add_rounded_rect(slide,
                         Inches(C2_LEFT + 1.2), Inches(BODY_TOP + 0.50),
                         Inches(2.0), Inches(1.6),
                         fill_color=theme["accent"], corner_size=0.15)
        add_textbox(slide,
                    Inches(C2_LEFT + 1.2), Inches(BODY_TOP + 1.0),
                    Inches(2.0), Inches(0.5),
                    "🌐 openclaw",
                    font_size=Pt(13), bold=True,
                    color=rgb(255,255,255), align=PP_ALIGN.CENTER)

    add_textbox(slide,
                Inches(C2_LEFT + PAD), Inches(BODY_TOP + 2.30),
                Inches(COL_W - PAD * 2), Inches(0.30),
                "openclaw browser",
                font_size=Pt(12), bold=True,
                color=theme["accent"], align=PP_ALIGN.CENTER)
    add_textbox(slide,
                Inches(C2_LEFT + PAD), Inches(BODY_TOP + 2.60),
                Inches(COL_W - PAD * 2), Inches(0.24),
                "--headless · 无头浏览器渲染",
                font_size=Pt(8.5),
                color=theme["body_text"], align=PP_ALIGN.CENTER)

    # VM grid: 2 rows × 3
    vm_y_start = BODY_TOP + 2.95
    vm_w = 1.10
    vm_h = 0.70
    vm_gap_x = 0.06
    vm_gap_y = 0.10
    vm_x_start = C2_LEFT + PAD + 0.12
    all_vms = VMS_EXISTING + VMS_EXTRA
    vm_colors = [theme["vm_fill"]] * 3 + [theme["vm_extra"]] * 3
    for idx, (vm_name, vm_color) in enumerate(zip(all_vms, vm_colors)):
        row = idx // 3
        col = idx % 3
        vx = vm_x_start + col * (vm_w + vm_gap_x)
        vy = vm_y_start + row * (vm_h + vm_gap_y)
        add_rounded_rect(slide,
                         Inches(vx), Inches(vy),
                         Inches(vm_w), Inches(vm_h),
                         fill_color=vm_color, corner_size=0.12)
        add_textbox(slide,
                    Inches(vx), Inches(vy + 0.05),
                    Inches(vm_w), Inches(0.24),
                    vm_name,
                    font_size=Pt(7.5), bold=True,
                    color=rgb(255,255,255), align=PP_ALIGN.CENTER)
        add_textbox(slide,
                    Inches(vx), Inches(vy + 0.28),
                    Inches(vm_w), Inches(0.20),
                    "2vCPU / 4G",
                    font_size=Pt(7),
                    color=rgb(220, 230, 245), align=PP_ALIGN.CENTER)
        add_textbox(slide,
                    Inches(vx), Inches(vy + 0.46),
                    Inches(vm_w), Inches(0.20),
                    "openclaw 🌐",
                    font_size=Pt(7),
                    color=rgb(255, 220, 180), align=PP_ALIGN.CENTER)

    # Overcommit ratio bar at bottom of col 2
    ratio_y = vm_y_start + 2 * (vm_h + vm_gap_y) + 0.12
    add_textbox(slide,
                Inches(C2_LEFT + PAD), Inches(ratio_y),
                Inches(COL_W - PAD * 2), Inches(0.24),
                "超分比  1.0×  →  ↑ 1.8×",
                font_size=Pt(10), bold=True,
                color=theme["accent"], align=PP_ALIGN.CENTER)

    # Arrow from col 2 → col 3
    add_textbox(slide,
                Inches(C2_LEFT + COL_W - 0.40), Inches(BODY_TOP + BODY_H / 2 - 0.15),
                Inches(0.55), Inches(0.35),
                "→",
                font_size=Pt(22), bold=True,
                color=theme["arrow_color"], align=PP_ALIGN.CENTER)

    # ------------------------------------------------------------------
    # 6. COLUMN 3 — SMAP memory layers
    # ------------------------------------------------------------------
    C3_LEFT = col_positions[2]
    layer_colors = [theme["l1_color"], theme["l2_color"], theme["l3_color"]]
    layer_icons  = ["dram_icon.svg", "ub_memory_icon.svg", "nvme_icon.svg"]
    layer_height = (BODY_H - 0.60) / 3   # divide body into 3 equal sections

    for li, (layer, lcolor, licon) in enumerate(
            zip(MEMORY_LAYERS, layer_colors, layer_icons)):
        ly_top = BODY_TOP + 0.48 + li * (layer_height + 0.04)

        # Layer panel
        add_rounded_rect(slide,
                         Inches(C3_LEFT + PAD), Inches(ly_top),
                         Inches(COL_W - PAD * 2 - 0.10), Inches(layer_height),
                         fill_color=theme["card_bg"],
                         line_color=lcolor,
                         line_width=Pt(1.5 if li == 2 else 0.75),
                         corner_size=0.12)

        # Color accent strip on left of panel
        add_rect(slide,
                 Inches(C3_LEFT + PAD), Inches(ly_top),
                 Inches(0.06), Inches(layer_height),
                 fill_color=lcolor)

        # Icon
        icon_png = svg_to_png_bytes(ASSETS_DIR / licon, 120, 120)
        icon_size = 0.80
        if icon_png:
            add_image_bytes(slide, icon_png,
                            Inches(C3_LEFT + PAD + 0.14), Inches(ly_top + (layer_height - icon_size) / 2),
                            Inches(icon_size), Inches(icon_size))
        else:
            # Fallback colored square
            add_rounded_rect(slide,
                             Inches(C3_LEFT + PAD + 0.14), Inches(ly_top + 0.10),
                             Inches(0.80), Inches(layer_height - 0.20),
                             fill_color=lcolor, corner_size=0.2)

        # Text block
        tx = C3_LEFT + PAD + 0.14 + icon_size + 0.10
        tw = COL_W - PAD * 2 - 0.10 - 0.14 - icon_size - 0.12

        add_textbox(slide,
                    Inches(tx), Inches(ly_top + 0.08),
                    Inches(tw), Inches(0.28),
                    layer["name"],
                    font_size=Pt(11), bold=True,
                    color=lcolor, align=PP_ALIGN.LEFT)
        add_textbox(slide,
                    Inches(tx), Inches(ly_top + 0.35),
                    Inches(tw), Inches(0.22),
                    f"{layer['sub']}  ·  {layer['data']}",
                    font_size=Pt(8.5),
                    color=theme["body_text"], align=PP_ALIGN.LEFT)
        add_textbox(slide,
                    Inches(tx), Inches(ly_top + 0.55),
                    Inches(tw), Inches(0.22),
                    f"访问延迟 {layer['latency']}",
                    font_size=Pt(8),
                    color=theme["body_text"], align=PP_ALIGN.LEFT)

        # L3 special badge "核心超分扩展"
        if li == 2:
            badge_x = C3_LEFT + PAD + 0.25
            badge_y = ly_top + layer_height - 0.35
            add_rounded_rect(slide,
                             Inches(badge_x), Inches(badge_y),
                             Inches(1.80), Inches(0.27),
                             fill_color=lcolor, corner_size=0.4)
            add_textbox(slide,
                        Inches(badge_x), Inches(badge_y),
                        Inches(1.80), Inches(0.27),
                        "★ 核心超分扩展",
                        font_size=Pt(8), bold=True,
                        color=rgb(255, 255, 255), align=PP_ALIGN.CENTER)

        # Down arrow between layers
        if li < 2:
            arrow_y = ly_top + layer_height + 0.00
            add_textbox(slide,
                        Inches(C3_LEFT + COL_W / 2 - 0.3), Inches(arrow_y),
                        Inches(0.6), Inches(0.09),
                        "↓ 冷数据降级",
                        font_size=Pt(7),
                        color=theme["body_text"], align=PP_ALIGN.CENTER)

    # ------------------------------------------------------------------
    # 7. Footer
    # ------------------------------------------------------------------
    footer_y = 7.5 - FOOTER_H
    add_rect(slide,
             Inches(0), Inches(footer_y), SLIDE_W, Inches(FOOTER_H),
             fill_color=theme["footer_bg"])

    add_textbox(slide,
                Inches(0.2), Inches(footer_y + 0.06),
                Inches(10), Inches(FOOTER_H - 0.08),
                "火山引擎 AI 平台  ·  openEuler OLK-6.6  ·  UB Stack  ·  SMAP 内存分层",
                font_size=Pt(8),
                color=theme["footer_text"], align=PP_ALIGN.LEFT)

    add_textbox(slide,
                Inches(11.0), Inches(footer_y + 0.06),
                Inches(2.2), Inches(FOOTER_H - 0.08),
                f"{theme['label']}",
                font_size=Pt(8),
                color=theme["footer_text"], align=PP_ALIGN.RIGHT)

    return prs


# ---------------------------------------------------------------------------
# Main
# ---------------------------------------------------------------------------
def main():
    print(f"Output directory: {OUTPUT_DIR}")
    print(f"SVG icons directory: {ASSETS_DIR}")
    print(f"cairosvg available: {HAS_CAIROSVG}")
    print()

    for i, theme in enumerate(THEMES, start=1):
        filename = f"ppt_beautiful_{theme['name']}.pptx"
        out_path = OUTPUT_DIR / filename

        print(f"[{i:02d}/10] Generating {filename} ({theme['label']}) ...", end="", flush=True)
        try:
            prs = build_slide(theme)
            prs.save(str(out_path))
            size_kb = out_path.stat().st_size // 1024
            print(f" ✓  ({size_kb} KB)")
        except Exception as e:
            print(f" ✗  ERROR: {e}")
            import traceback
            traceback.print_exc()

    print()
    print("Done! Generated files:")
    for f in sorted(OUTPUT_DIR.glob("ppt_beautiful_*.pptx")):
        print(f"  {f.name}  ({f.stat().st_size // 1024} KB)")


if __name__ == "__main__":
    main()
